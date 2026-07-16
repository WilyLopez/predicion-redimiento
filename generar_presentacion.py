from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os


def crear_presentacion_profesional():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ─────────────────────────────────────────────
    #  PALETA DE COLORES — Dark Premium Theme
    # ─────────────────────────────────────────────
    C_BG       = RGBColor(8,   12,  24)    # Fondo oscuro (casi negro navy)
    C_BG2      = RGBColor(15,  23,  42)    # Slate 900
    C_CARD     = RGBColor(20,  30,  50)    # Tarjeta oscura
    C_BORDER   = RGBColor(44,  60,  95)    # Borde sutil
    C_BLUE     = RGBColor(59,  130, 246)   # Blue 500
    C_BLUE_D   = RGBColor(37,  99,  235)   # Blue 600 (acento)
    C_BLUE_L   = RGBColor(96,  165, 250)   # Blue 400 (claro)
    C_GREEN    = RGBColor(16,  185, 129)   # Emerald 500
    C_GREEN_D  = RGBColor(5,   150, 105)   # Emerald 600
    C_AMBER    = RGBColor(245, 158, 11)    # Amber 500
    C_RED      = RGBColor(239, 68,  68)    # Red 500
    C_RED_D    = RGBColor(220, 38,  38)    # Red 600
    C_WHITE    = RGBColor(248, 250, 252)   # Casi blanco
    C_GRAY1    = RGBColor(148, 163, 184)   # Slate 400 (texto secundario)
    C_GRAY2    = RGBColor(100, 116, 139)   # Slate 500 (texto apagado)
    C_PURPLE   = RGBColor(139, 92,  246)   # Violet 500
    C_TEAL     = RGBColor(20,  184, 166)   # Teal 500
    C_BLACK    = RGBColor(4,   6,   12)    # Sombra

    # ─────────────────────────────────────────────
    #  HELPERS
    # ─────────────────────────────────────────────

    def add_bg(slide, color=C_BG2):
        """Fondo completo de la diapositiva."""
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()

    def add_left_stripe(slide, color=C_BLUE_D, w=0.45):
        """Franja vertical izquierda de acento."""
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(w), prs.slide_height)
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
        # Segunda franja más delgada
        b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(w), 0, Inches(0.08), prs.slide_height)
        b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(30, 80, 190); b2.line.fill.background()

    def add_footer(slide, texto="Analisis de Riesgo Estudiantil  |  Python · XGBoost · SHAP · Streamlit"):
        """Barra de pie de página en cada diapositiva."""
        fh = Inches(0.28)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - fh, prs.slide_width, fh)
        bar.fill.solid(); bar.fill.fore_color.rgb = C_BLUE_D; bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - fh, Inches(12.5), fh)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = texto; p.font.name = "Arial"; p.font.size = Pt(9)
        p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

    def add_deco_circle(slide, x, y, d, color):
        """Círculo decorativo de fondo."""
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()

    def add_slide_title(slide, text, x=0.65, y=0.18, size=30):
        """Título estándar con subrayado de acento."""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(12.3), Inches(0.75))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text; p.font.name = "Arial"; p.font.size = Pt(size)
        p.font.bold = True; p.font.color.rgb = C_WHITE
        # Línea subrayado
        ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.62), Inches(4.0), Inches(0.04))
        ul.fill.solid(); ul.fill.fore_color.rgb = C_BLUE; ul.line.fill.background()

    def add_card(slide, x, y, w, h, color=None, border=None, rounded=True):
        """Tarjeta con sombra debajo."""
        st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        # Sombra
        sh = slide.shapes.add_shape(st, Inches(x + 0.07), Inches(y + 0.07), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = C_BLACK; sh.line.fill.background()
        # Tarjeta principal
        cd = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
        cd.fill.solid(); cd.fill.fore_color.rgb = color or C_CARD
        if border:
            cd.line.color.rgb = border; cd.line.width = Pt(1.5)
        else:
            cd.line.fill.background()
        return cd

    def add_top_accent(slide, x, y, w, h=0.11, color=C_BLUE):
        """Franja de color en el borde superior de una tarjeta."""
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()

    def add_left_indicator(slide, x, y, h, color):
        """Indicador vertical de color en el borde izquierdo de una tarjeta."""
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()

    def add_text(slide, x, y, w, h, text, size=13, color=None,
                 bold=False, align=PP_ALIGN.LEFT, italic=False):
        """Caja de texto rápida."""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text; p.font.name = "Arial"; p.font.size = Pt(size)
        p.font.bold = bold; p.font.italic = italic
        p.font.color.rgb = color or C_WHITE; p.alignment = align
        return tb

    def add_progress_bar(slide, x, y, w, h, pct, color=C_BLUE):
        """Barra de progreso visual."""
        bg_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg_b.fill.solid(); bg_b.fill.fore_color.rgb = RGBColor(30, 41, 59); bg_b.line.fill.background()
        fw = max(0.01, w * (pct / 100.0))
        fg_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(fw), Inches(h))
        fg_b.fill.solid(); fg_b.fill.fore_color.rgb = color; fg_b.line.fill.background()

    def add_badge_circle(slide, x, y, d, color, text, text_size=14):
        """Círculo badge numerado."""
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
        add_text(slide, x, y + d * 0.1, d, d * 0.8, text, size=text_size, bold=True, align=PP_ALIGN.CENTER)

    def add_pill_badge(slide, x, y, w, h, text, color, text_size=10):
        """Pequeño badge tipo píldora."""
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()
        add_text(slide, x, y + 0.02, w, h, text, size=text_size, bold=True, align=PP_ALIGN.CENTER)

    def add_divider(slide, x, y, w, color=C_BORDER):
        d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.03))
        d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 1 — PORTADA
    # ══════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1, C_BG)

    # Círculos decorativos de fondo (top-right)
    add_deco_circle(s1, 10.2,  -2.0, 6.0, RGBColor(18, 35,  90))
    add_deco_circle(s1, 11.0,  -1.0, 4.0, RGBColor(28, 55, 140))
    add_deco_circle(s1, 11.8,  -0.2, 2.2, RGBColor(37, 80, 200))
    # Círculo decorativo bottom-left
    add_deco_circle(s1, -1.5,   5.2, 3.5, RGBColor(12, 25,  65))

    # Franja izquierda
    s1_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.65), prs.slide_height)
    s1_bar.fill.solid(); s1_bar.fill.fore_color.rgb = C_BLUE_D; s1_bar.line.fill.background()
    s1_bar2 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), 0, Inches(0.1), prs.slide_height)
    s1_bar2.fill.solid(); s1_bar2.fill.fore_color.rgb = RGBColor(30, 75, 185); s1_bar2.line.fill.background()

    # Título principal
    tb_main = s1.shapes.add_textbox(Inches(1.1), Inches(1.0), Inches(10.8), Inches(2.4))
    tf_main = tb_main.text_frame; tf_main.word_wrap = True
    p1 = tf_main.paragraphs[0]
    p1.text = "Analisis de Riesgo y Prediccion"
    p1.font.name = "Arial"; p1.font.size = Pt(48); p1.font.bold = True; p1.font.color.rgb = C_WHITE
    p2 = tf_main.add_paragraph()
    p2.text = "del Rendimiento Estudiantil"
    p2.font.name = "Arial"; p2.font.size = Pt(48); p2.font.bold = True; p2.font.color.rgb = C_BLUE_L

    # Subtítulo
    add_text(s1, 1.1, 3.55, 10.5, 0.55,
             "Deteccion temprana de la desercion universitaria mediante Machine Learning e IA Explicable",
             size=18, color=C_GRAY1)

    # Línea divisora horizontal
    add_divider(s1, 1.1, 4.3, 5.5, C_BLUE_D)

    # Tags tecnología
    tags = [("Python", C_GREEN_D), ("XGBoost", C_BLUE_D), ("SHAP", C_PURPLE), ("Streamlit", C_RED_D), ("K-Means", RGBColor(15, 150, 155))]
    tx = 1.1
    for tlabel, tcolor in tags:
        add_pill_badge(s1, tx, 4.55, 1.55, 0.42, tlabel, tcolor, text_size=13)
        tx += 1.72

    # Barra de estadísticas inferior
    stats_bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.1), prs.slide_width, Inches(0.9))
    stats_bg.fill.solid(); stats_bg.fill.fore_color.rgb = RGBColor(12, 20, 38); stats_bg.line.fill.background()

    stats = [("76.6%", "Accuracy XGBoost"), ("32.1%", "Tasa Desercion"), ("3", "Paradigmas"), ("3", "Algoritmos ML"), ("13", "Pruebas Unitarias")]
    sx = 0.8
    for sv, sl in stats:
        add_text(s1, sx, 6.13, 2.3, 0.45, sv, size=24, bold=True, color=C_BLUE_L, align=PP_ALIGN.CENTER)
        add_text(s1, sx, 6.57, 2.3, 0.28, sl, size=10, color=C_GRAY2, align=PP_ALIGN.CENTER)
        # Separador vertical
        if sx < 9.5:
            sep = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx + 2.3), Inches(6.18), Inches(0.02), Inches(0.6))
            sep.fill.solid(); sep.fill.fore_color.rgb = C_BORDER; sep.line.fill.background()
        sx += 2.5

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 2 — INTRODUCCIÓN
    # ══════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2); add_left_stripe(s2); add_footer(s2)
    add_deco_circle(s2, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s2, "Introduccion y Contexto de la Desercion", x=0.7, y=0.18)

    # Tarjeta stat grande (izquierda) — 32.1%
    add_card(s2, 0.7, 1.05, 4.65, 5.65, color=RGBColor(38, 8, 8), border=C_RED_D)
    add_top_accent(s2, 0.7, 1.05, 4.65, 0.13, C_RED)

    add_text(s2, 0.85, 1.3, 4.35, 1.5, "32.1%", size=82, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    add_text(s2, 0.85, 2.9, 4.35, 0.45, "Tasa Historica de Desercion", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_divider(s2, 1.2, 3.47, 3.6, RGBColor(80, 18, 18))
    add_text(s2, 0.9, 3.6, 4.25, 2.8,
             "La desercion en educacion superior es un problema multifactorial que afecta a estudiantes e instituciones. "
             "La deteccion y alerta temprana durante el primer semestre es vital para mejorar la retencion.",
             size=13, color=C_GRAY1, align=PP_ALIGN.CENTER)

    # Tarjeta 1 — El Desafío (derecha arriba)
    add_card(s2, 5.7, 1.05, 7.25, 2.6, color=C_CARD, border=C_BORDER)
    add_top_accent(s2, 5.7, 1.05, 7.25, 0.12, C_AMBER)
    add_text(s2, 5.9, 1.23, 6.85, 0.42, "El Desafio de la Intervencion Academica", size=15, bold=True, color=C_AMBER)
    add_divider(s2, 5.9, 1.72, 6.8, RGBColor(60, 45, 15))
    add_text(s2, 5.9, 1.85, 6.85, 1.65,
             "Las instituciones actuan tarde: cuando el alumno ya ha reprobado o abandonado. "
             "Identificar alertas financieras y academicas desde el primer semestre permite una respuesta agil y focalizada.",
             size=13, color=C_GRAY1)

    # Tarjeta 2 — La Solución (derecha abajo)
    add_card(s2, 5.7, 3.9, 7.25, 2.6, color=C_CARD, border=C_BORDER)
    add_top_accent(s2, 5.7, 3.9, 7.25, 0.12, C_GREEN)
    add_text(s2, 5.9, 4.08, 6.85, 0.42, "La Solucion Propuesta", size=15, bold=True, color=C_GREEN)
    add_divider(s2, 5.9, 4.57, 6.8, RGBColor(12, 50, 35))
    add_text(s2, 5.9, 4.7, 6.85, 1.65,
             "Un sistema integral que calcula el score de riesgo exacto de cada estudiante, "
             "identifica su perfil socioeconomico mediante clustering y explica los motivos con Inteligencia Artificial Explicable (SHAP).",
             size=13, color=C_GRAY1)

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 3 — DATASET
    # ══════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3); add_left_stripe(s3); add_footer(s3)
    add_deco_circle(s3, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s3, "Estructura de Datos del Estudiante", x=0.7, y=0.18)

    cols = [
        {"icon": "◉", "titulo": "Variables Demograficas", "color": C_BLUE, "x": 0.7,
         "items": ["Edad al inscribirse (factor critico).", "Genero e informacion de desplazamiento.", "Estado civil y nacionalidad.", "Condicion de estudiante internacional."]},
        {"icon": "$", "titulo": "Factores Socioeconomicos", "color": C_AMBER, "x": 4.97,
         "items": ["Becario (principal factor protector).", "Deudor de matricula y pagos al dia.", "Instruccion y ocupacion de los padres.", "Necesidades educativas especiales."]},
        {"icon": "★", "titulo": "Desempeno Academico", "color": C_GREEN, "x": 9.24,
         "items": ["U. curriculares aprobadas sem. 1 y 2.", "Calificaciones promedio iniciales.", "Unidades homologadas o acreditadas.", "Nota de admision al programa."]},
    ]

    for col in cols:
        add_card(s3, col["x"], 1.05, 3.92, 5.7, color=C_CARD, border=C_BORDER)
        add_top_accent(s3, col["x"], 1.05, 3.92, 0.12, col["color"])

        # Badge con ícono
        add_badge_circle(s3, col["x"] + 0.15, 1.22, 0.55, col["color"], col["icon"], text_size=16)

        add_text(s3, col["x"] + 0.2, 1.88, 3.52, 0.44, col["titulo"], size=15, bold=True, color=C_WHITE)
        add_divider(s3, col["x"] + 0.2, 2.37, 3.35, C_BORDER)

        iy = 2.52
        for item in col["items"]:
            dot = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(col["x"] + 0.2), Inches(iy + 0.1), Inches(0.1), Inches(0.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = col["color"]; dot.line.fill.background()
            add_text(s3, col["x"] + 0.38, iy, 3.38, 0.58, item, size=12, color=C_GRAY1)
            iy += 0.65

    # Banner inferior del dataset
    add_card(s3, 0.7, 6.8, 12.25, 0.42, color=RGBColor(18, 30, 60), border=C_BLUE_D, rounded=False)
    add_text(s3, 1.0, 6.85, 11.7, 0.32,
             "Dataset UCI  ·  ~4,000 estudiantes  ·  36 variables de entrada  ·  Objetivo: Dropout / Enrolled / Graduate",
             size=12, color=C_GRAY1, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 4 — ARQUITECTURA
    # ══════════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4); add_left_stripe(s4); add_footer(s4)
    add_deco_circle(s4, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s4, "Arquitectura y Flujo del Software", x=0.7, y=0.18)

    steps = [
        {"n": "1", "name": "Datos CSV", "desc": "Lectura con autodetector de delimitador", "color": C_BLUE_L,  "x": 0.7},
        {"n": "2", "name": "Procesador", "desc": "Limpieza, normalizacion y split train/test", "color": C_BLUE,   "x": 3.25},
        {"n": "3", "name": "Clasificador", "desc": "XGBoost, Random Forest, Reg. Logistica", "color": C_PURPLE, "x": 5.8},
        {"n": "4", "name": "Explicabilidad", "desc": "Valores de impacto SHAP por variable", "color": C_AMBER,  "x": 8.35},
        {"n": "5", "name": "Dashboard", "desc": "Interfaz interactiva con Streamlit", "color": C_GREEN,  "x": 10.9},
    ]

    for i, step in enumerate(steps):
        add_card(s4, step["x"], 1.05, 2.17, 2.25, color=C_CARD, border=step["color"])
        add_top_accent(s4, step["x"], 1.05, 2.17, 0.12, step["color"])
        add_badge_circle(s4, step["x"] + 0.86, 1.22, 0.45, step["color"], step["n"], text_size=14)
        add_text(s4, step["x"] + 0.1, 1.78, 1.97, 0.42, step["name"], size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s4, step["x"] + 0.1, 2.25, 1.97, 0.9, step["desc"], size=11, color=C_GRAY1, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arr = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(step["x"] + 2.21), Inches(2.05), Inches(0.08), Inches(0.32))
            arr.fill.solid(); arr.fill.fore_color.rgb = step["color"]; arr.line.fill.background()

    # Paradigmas
    paras = [
        {"label": "POO — Orientado a Objetos", "color": C_BLUE,
         "desc": "Clases independientes: ProcesadorDatos, ClasificadorRendimiento, AgrupadorEstudiantes y ExplicadorRiesgo. Encapsulan el estado del pipeline para reusabilidad limpia."},
        {"label": "PF — Programacion Funcional", "color": C_GREEN,
         "desc": "Decoradores @medir_tiempo y @capturar_errores. Extienden metodos sin modificar su logica interna. Generan logs de auditoria y tiempo automaticamente."},
        {"label": "PP — Procedural", "color": C_AMBER,
         "desc": "Scripts de orquestacion lineal: descargar_datos.py y test_agrupamiento_explicabilidad.py. Flujo secuencial de descarga, validacion y ejecucion de pruebas."},
    ]
    px = 0.7
    for para in paras:
        add_card(s4, px, 3.55, 4.08, 2.85, color=RGBColor(14, 22, 40), border=para["color"])
        add_left_indicator(s4, px, 3.55, 2.85, para["color"])
        add_text(s4, px + 0.27, 3.68, 3.65, 0.42, para["label"], size=13, bold=True, color=para["color"])
        add_divider(s4, px + 0.27, 4.15, 3.55, C_BORDER)
        add_text(s4, px + 0.27, 4.25, 3.65, 2.0, para["desc"], size=12, color=C_GRAY1)
        px += 4.3

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 5 — MODELOS
    # ══════════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5); add_left_stripe(s5); add_footer(s5)
    add_deco_circle(s5, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s5, "Resultados de Modelos Predictivos", x=0.7, y=0.18)

    # Tabla oscura
    t = s5.shapes.add_table(4, 4, Inches(0.7), Inches(1.05), Inches(7.8), Inches(3.2)).table
    t.columns[0].width = Inches(1.9); t.columns[1].width = Inches(1.5)
    t.columns[2].width = Inches(2.5); t.columns[3].width = Inches(1.9)

    headers = ["Algoritmo", "Accuracy", "Ventaja Principal", "Limitacion"]
    rows_data = [
        ["XGBoost  ✓", "76.61%", "Maxima capacidad predictiva. Elegido como modelo principal.", "Mayor costo computacional con SHAP."],
        ["Random Forest", "75.03%", "Muy estable. Rapido en calculo de SHAP.", "Menor precision ante relaciones no lineales."],
        ["Reg. Logistica", "74.01%", "Extrema interpretabilidad. Linea base estadistica.", "Asume linealidad entre los factores."],
    ]

    for ci, h in enumerate(headers):
        cell = t.cell(0, ci); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = C_BLUE_D
        for p in cell.text_frame.paragraphs:
            p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = C_WHITE; p.font.name = "Arial"

    row_bg = [RGBColor(20, 35, 60), RGBColor(16, 26, 46)]
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = t.cell(ri + 1, ci); cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = row_bg[ri % 2]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12); p.font.name = "Arial"
                p.font.color.rgb = (C_GREEN if (ri == 0 and ci <= 1) else C_GRAY1)
                p.font.bold = (ri == 0 and ci <= 1)

    # Panel derecho — Barras visuales de accuracy
    add_card(s5, 8.8, 1.05, 4.15, 3.2, color=RGBColor(14, 22, 40), border=C_BORDER)
    add_text(s5, 8.95, 1.1, 3.85, 0.4, "Comparativa Visual (Accuracy)", size=13, bold=True, color=C_GRAY1)

    models = [("XGBoost", 76.61, C_GREEN), ("Random Forest", 75.03, C_BLUE), ("Reg. Logistica", 74.01, C_AMBER)]
    by = 1.62
    for mname, macc, mcol in models:
        add_text(s5, 8.95, by, 2.3, 0.32, mname, size=12, color=C_GRAY1)
        add_text(s5, 11.55, by, 1.25, 0.32, f"{macc}%", size=12, bold=True, color=mcol, align=PP_ALIGN.RIGHT)
        add_progress_bar(s5, 8.95, by + 0.35, 3.82, 0.25, macc, mcol)
        by += 0.95

    # F1-Score strip
    add_card(s5, 0.7, 4.38, 12.25, 1.85, color=RGBColor(10, 18, 40), border=C_BORDER, rounded=False)
    add_text(s5, 0.9, 4.48, 12.0, 0.38, "F1-Score por Clase — XGBoost (Modelo Principal)", size=13, bold=True, color=C_BLUE_L)
    add_divider(s5, 0.9, 4.9, 11.8, C_BORDER)

    f1s = [
        ("Desercion  (Dropout)", "F1 = 0.77", "Detecta 3 de cada 4 alumnos en riesgo real.", C_RED),
        ("Graduado  (Graduate)", "F1 = 0.84", "La clase mas predecible. Notas altas y pagos al dia.", C_GREEN),
        ("Matriculado  (Enrolled)", "F1 = 0.49", "Transicion inestable: no deserto pero tampoco graduo.", C_AMBER),
    ]
    fx = 0.9
    for fn, fv, fd, fc in f1s:
        add_text(s5, fx, 5.0, 3.95, 0.35, fn, size=12, bold=True, color=fc)
        add_text(s5, fx, 5.38, 3.95, 0.7, fv + "  —  " + fd, size=11, color=C_GRAY1)
        fx += 4.12

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 6 — CLUSTERING + SHAP
    # ══════════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s6); add_left_stripe(s6); add_footer(s6)
    add_deco_circle(s6, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s6, "Clustering Socioeconomico y Explicabilidad (SHAP)", x=0.7, y=0.18)

    # ── Columna izquierda: Clustering ──
    add_card(s6, 0.7, 1.05, 6.1, 5.75, color=C_CARD, border=C_BORDER)
    add_top_accent(s6, 0.7, 1.05, 6.1, 0.12, C_BLUE)
    add_text(s6, 0.9, 1.22, 5.7, 0.42, "K-Means — 3 Perfiles de Vulnerabilidad", size=16, bold=True, color=C_BLUE_L)
    add_text(s6, 0.9, 1.68, 5.7, 0.32, "Silhouette Score calcula y valida automaticamente la eleccion de K=3", size=11, color=C_GRAY2, italic=True)
    add_divider(s6, 0.9, 2.07, 5.6, C_BORDER)

    perfiles = [
        {"id": "0", "titulo": "Becados y Estables", "stats": "Edad ~19.5 anos  |  Becarios: 85-90%  |  Deudores: 2%", "nivel": "RIESGO BAJO",  "nc": C_GREEN, "color": C_GREEN},
        {"id": "1", "titulo": "Mayores en Alta Vulnerabilidad", "stats": "Edad ~29 anos  |  Becarios: <5%  |  Deudores: 35-40%", "nivel": "RIESGO ALTO",  "nc": C_RED,   "color": C_RED},
        {"id": "2", "titulo": "Jovenes sin Beca", "stats": "Edad ~20.2 anos  |  Becarios: 0%  |  Deudores: 12%", "nivel": "RIESGO MEDIO", "nc": C_AMBER, "color": C_AMBER},
    ]
    py = 2.22
    for perf in perfiles:
        add_card(s6, 0.85, py, 5.72, 1.08, color=RGBColor(14, 24, 46), border=perf["color"])
        add_left_indicator(s6, 0.85, py, 1.08, perf["color"])
        add_badge_circle(s6, 1.1, py + 0.32, 0.38, perf["color"], perf["id"], text_size=14)
        add_text(s6, 1.62, py + 0.1, 3.3, 0.4, perf["titulo"], size=13, bold=True, color=C_WHITE)
        add_text(s6, 1.62, py + 0.55, 3.3, 0.35, perf["stats"], size=11, color=C_GRAY1)
        add_pill_badge(s6, 5.07, py + 0.32, 1.37, 0.37, perf["nivel"], perf["nc"], text_size=9)
        py += 1.22

    add_text(s6, 0.9, 5.72, 5.7, 0.75,
             "El clustering corre el dia de la matricula, ANTES del primer examen. Identifica el grupo de riesgo sin necesidad de calificaciones.",
             size=12, color=C_GRAY2, italic=True)

    # ── Columna derecha: SHAP ──
    add_card(s6, 7.15, 1.05, 5.8, 5.75, color=C_CARD, border=C_BORDER)
    add_top_accent(s6, 7.15, 1.05, 5.8, 0.12, C_PURPLE)
    add_text(s6, 7.35, 1.22, 5.4, 0.42, "Auditoria de Riesgo con SHAP (XAI)", size=16, bold=True, color=RGBColor(167, 139, 250))
    add_divider(s6, 7.35, 1.7, 5.3, C_BORDER)

    shap_items = [
        ("Apertura de Caja Negra", "Explica que factores aumentan o reducen la desercion de cada alumno de forma individual, no como promedio.", C_BLUE),
        ("Accion Localizable", "El tutor identifica si el riesgo es financiero (deudas) o academico (notas bajas) para derivar el apoyo correcto.", C_GREEN),
        ("Grafica Dinamica", "Barras rojas: factores de riesgo (deudas, reprobacion). Barras azules: factores protectores (beca, materias aprobadas).", C_AMBER),
        ("Importancia Global", "Calcula el impacto medio |SHAP| sobre la poblacion completa para identificar las variables mas criticas del sistema.", C_PURPLE),
    ]
    sy = 1.82
    for st_, sd, sc in shap_items:
        dot = s6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), Inches(sy + 0.12), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = sc; dot.line.fill.background()
        add_text(s6, 7.55, sy, 5.2, 0.38, st_, size=13, bold=True, color=sc)
        add_text(s6, 7.55, sy + 0.4, 5.2, 0.78, sd, size=12, color=C_GRAY1)
        sy += 1.3

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 7 — CONCLUSIONES
    # ══════════════════════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s7); add_left_stripe(s7); add_footer(s7)
    add_deco_circle(s7, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s7, "Conclusiones e Impacto del Sistema", x=0.7, y=0.18)

    # Tarjeta resumen izquierda
    add_card(s7, 0.7, 1.05, 6.1, 2.65, color=C_CARD, border=C_BORDER)
    add_top_accent(s7, 0.7, 1.05, 6.1, 0.12, C_BLUE)
    add_text(s7, 0.9, 1.22, 5.7, 0.42, "Prevencion Inteligente Basada en Datos", size=15, bold=True, color=C_BLUE_L)
    add_divider(s7, 0.9, 1.7, 5.7, RGBColor(25, 45, 85))
    add_text(s7, 0.9, 1.82, 5.7, 1.75,
             "El sistema reemplaza auditorias reactivas por predicciones oportunas con 76.6% de accuracy. "
             "Los tutores focalizan becas y tutorias en el grupo prioritario, optimizando los recursos institucionales.",
             size=13, color=C_GRAY1)

    # Tarjeta resumen derecha
    add_card(s7, 7.15, 1.05, 5.8, 2.65, color=C_CARD, border=C_BORDER)
    add_top_accent(s7, 7.15, 1.05, 5.8, 0.12, C_GREEN)
    add_text(s7, 7.35, 1.22, 5.4, 0.42, "Valor Academico y Tecnico", size=15, bold=True, color=C_GREEN)
    add_divider(s7, 7.35, 1.7, 5.3, RGBColor(10, 45, 30))
    add_text(s7, 7.35, 1.82, 5.4, 1.75,
             "Arquitectura modular en Python bajo 3 paradigmas. Incluye manual de usuario, reporte de modelos, "
             "13 pruebas unitarias y dashboard listo para produccion. Cumple integramente la rubrica de excelencia.",
             size=13, color=C_GRAY1)

    # 5 Badges de logros
    ach = [
        ("76.6%", "Accuracy\nXGBoost", C_GREEN),
        ("F1: 0.77", "Deteccion de\nDesercion", C_BLUE),
        ("Dia 1", "Clustering sin\ncalificaciones", C_AMBER),
        ("SHAP", "XAI por\nalumno", C_PURPLE),
        ("3", "Paradigmas\nde codigo", C_TEAL),
    ]
    ax = 0.7
    for av, al, ac in ach:
        add_card(s7, ax, 3.9, 2.42, 2.12, color=RGBColor(13, 22, 40), border=ac)
        add_top_accent(s7, ax, 3.9, 2.42, 0.12, ac)
        add_text(s7, ax + 0.05, 4.07, 2.32, 0.72, av, size=28, bold=True, color=ac, align=PP_ALIGN.CENTER)
        add_divider(s7, ax + 0.2, 4.82, 2.02, RGBColor(35, 50, 80))
        add_text(s7, ax + 0.05, 4.9, 2.32, 1.0, al, size=12, color=C_GRAY1, align=PP_ALIGN.CENTER)
        ax += 2.54

    # Frase de cierre
    add_card(s7, 0.7, 6.17, 12.25, 0.65, color=C_BLUE_D, border=C_BLUE, rounded=False)
    add_text(s7, 1.0, 6.24, 11.8, 0.52,
             "\"Este sistema no reemplaza al tutor academico, lo potencia. La informacion correcta, en el momento correcto, sobre el alumno correcto.\"",
             size=13, bold=True, align=PP_ALIGN.CENTER, color=C_WHITE)

    # ══════════════════════════════════════════════════════════════════════
    #  DIAPOSITIVA 8 — Q&A
    # ══════════════════════════════════════════════════════════════════════
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s8); add_left_stripe(s8); add_footer(s8)
    add_deco_circle(s8, 10.5, -1.0, 3.5, RGBColor(14, 28, 72))
    add_slide_title(s8, "Preguntas Frecuentes del Jurado", x=0.7, y=0.18)

    qa = [
        {"q": "P1:  ¿Por que XGBoost y no una red neuronal?",
         "a": "Con ~4,000 registros, una red neuronal no justifica su complejidad. XGBoost maximiza el accuracy con menor costo computacional "
              "y es nativamente compatible con SHAP para generar explicaciones locales por alumno.",
         "color": C_BLUE},
        {"q": "P2:  ¿El modelo no discrimina por genero o etnia?",
         "a": "Segun el analisis SHAP global, los factores academicos y financieros dominan las predicciones. El genero es variable secundaria. "
              "Se usa class_weight='balanced' para evitar sesgos de clase en el entrenamiento.",
         "color": C_AMBER},
        {"q": "P3:  ¿Como se justifica K=3 en el clustering?",
         "a": "Se calcula el Silhouette Score para K=2 hasta K=6 (visible en la pestana 'Validacion del Modelo' del dashboard). "
              "K=3 maximiza la separacion entre clusters en los datos socioeconomicos multivariados del dataset.",
         "color": C_GREEN},
    ]

    qy = 1.05
    for q_item in qa:
        add_card(s8, 0.7, qy, 12.25, 1.75, color=C_CARD, border=q_item["color"])
        add_left_indicator(s8, 0.7, qy, 1.75, q_item["color"])
        add_text(s8, 0.98, qy + 0.1, 11.8, 0.44, q_item["q"], size=14, bold=True, color=q_item["color"])
        add_divider(s8, 0.98, qy + 0.58, 11.7, C_BORDER)
        add_text(s8, 0.98, qy + 0.68, 11.8, 0.95, q_item["a"], size=13, color=C_GRAY1)
        qy += 1.95

    # ══════════════════════════════════════════════════════════════════════
    #  GUARDAR
    # ══════════════════════════════════════════════════════════════════════
    nombre = "presentacion_proyecto.pptx"
    try:
        prs.save(nombre)
        print(f"[OK] Presentacion premium generada: {nombre}  ({prs.slides.__len__()} diapositivas)")
    except PermissionError:
        alt = "presentacion_proyecto_v2.pptx"
        prs.save(alt)
        print(f"[OK] Archivo en uso, guardado como: {alt}")


if __name__ == "__main__":
    crear_presentacion_profesional()
